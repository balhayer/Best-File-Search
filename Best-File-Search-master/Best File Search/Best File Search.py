import os
import tkinter
from tkinter import *
from PyPDF2 import PdfFileReader #PDF
import comtypes.client #DOCX
import win32api # install pypiwin32, PyMuPDF, python-pptx, python-docx
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from pptx import Presentation
import xlsxwriter #EXCEL
from xlsxwriter.utility import xl_rowcol_to_cell
import xlrd
import fitz # PDF


def init(window):
    window.title("Best File Search")
    window.geometry("700x640")
    window.resizable(False, False)

    # Configure Widgets
    frame.place(x=300, y=78)  # set for fileList(Listbox) scrollbar
    frame2.place(x=300, y=412)
    scrollbar.configure(command=fileList.yview)
    fileList.configure(yscrollcommand=scrollbar.set)
    scrollbar2.configure(command=fileList.yview)
    openedFileList.configure(yscrollcommand=scrollbar.set)
    backgroundlabel.configure(image=backgroundimage)
    bClear.configure(command=buttonClear)
    bOpen.configure(command=buttonOpen)
    btnSearch.configure(command=executeSearchButton)
    bClearTwo.configure(command=buttonClearTwo)
    # Place Widgets
    # labelsearchtext.place(x=40, y=80)
    # labledesiredpath.place(x=40, y=230)
    desiredpath.place(x=130, y=250)
    searchText.place(x=41, y=125)
    # searchedFileName.place(x=300, y=30)
    # openedFileName.place(x=300, y=380)
    btnSearch.place(x=170, y=110, width=50, height=50)
    bClear.place(x=620, y=385, width=50, height=20)
    bClearTwo.place(x=565,y=50, width=50, height=20)
    bOpen.place(x=620, y=50, width=50, height=20)

def buttonClear():
    openedFileList.delete(0, END)

def buttonClearTwo():
    fileList.delete(0, END)

def buttonOpen():
    value = fileList.get(fileList.curselection()[0])
    #os.startfile(value)

    searchInput = searchText.get()

    ext = os.path.splitext(value)[-1]
    pathExceptExt = os.path.splitext(value)[0]

    if ext == ".pptx":
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        Filename = value
        deck = powerpoint.Presentations.Open(Filename)
        deck.SaveAs(pathExceptExt + ".pdf", 32)  # formatType = 32 for ppt to pdf

        value2 = pathExceptExt + ".pdf"
        doc = fitz.open(value2)
        pdf = PdfFileReader(value2)
        saveNumPages = pdf.getNumPages()

        for i in range(saveNumPages):
            page = doc[i]
            text_instances = page.searchFor(searchInput)

            for inst in text_instances:
                highlight = page.addHighlightAnnot(inst)
        doc.save(pathExceptExt + "_modified.pdf", garbage=4, deflate=True, clean=True)
        os.startfile(pathExceptExt + "_modified.pdf")

    elif ext == ".txt":
        with open(value, mode='r', encoding='utf-8') as fo:
            textthis = fo.read()

        with open(pathExceptExt + ".html", "w") as e:
            if searchInput in textthis:
                e.write("<pre>" + textthis.replace(searchInput, '<span style="background-color: #FFFF00">{}</span>'.format(
                    searchInput)) + "</pre> <br>\n")
        os.startfile(pathExceptExt + ".html")
    elif ext == ".xlsx":

        wbk = xlsxwriter.Workbook(value)
        wks = wbk.add_worksheet()
        myPath = value

        cell_format = wbk.add_format()
        cell_format.set_bg_color('yellow')
        for sh in xlrd.open_workbook(myPath).sheets():
            for row in range(sh.nrows):
                for col in range(sh.ncols):
                    mycell = sh.cell(row, col)
                    if mycell.value == searchInput:
                        wks.write(xl_rowcol_to_cell(row, col), searchInput, cell_format)
                    else:
                        wks.write(xl_rowcol_to_cell(row, col), mycell.value)
        wbk.close()
        os.startfile(value)
    elif ext == ".docx":
        document = Document(value)

        for para in document.paragraphs:
            start = para.text.find(searchInput)
            if start > -1:
                pre = para.text[:start]
                post = para.text[start + len(searchInput):]
                para.text = pre
                para.add_run(searchInput)
                para.runs[1].font.highlight_color = WD_COLOR_INDEX.YELLOW
                para.add_run(post)
        document.save(pathExceptExt + "_modified.docx")
        os.startfile(pathExceptExt + "_modified.docx")

    elif ext == ".pdf":

        doc = fitz.open(value)
        for i in range(10):
            try:
                    page = doc[i]

                    text_instances = page.searchFor(searchInput)

                    for inst in text_instances:
                        highlight = page.addHighlightAnnot(inst)

            except IndexError:
                break
        doc.save(pathExceptExt + "_mod.pdf", garbage=4, deflate=True, clean=True)
        os.startfile(pathExceptExt + "_mod.pdf")

    openedFileList.insert(END, value)

# Determine whether there are specific strings in some files such as .txt, .docx, .pptx files
# then insert their path and name to the list
def searchDir(root_folder, searchInput):

    #exceptFolderList = [ "AppData", "WINDOWS", "Windows", "Program Files (x86)", "Program Files"]

    try:
        filenames = os.listdir(root_folder)

        for filename in filenames:
            full_filename = os.path.join(root_folder, filename)

            if os.path.isdir(full_filename):
                if filename == "AppData" or filename == "WINDOWS" or filename == "Windows" or filename == "Program Files (x86)" or filename == "Program Files":
                    pass
                else:
                    searchDir(full_filename, searchInput)

            else:
                ext = os.path.splitext(full_filename)[-1]
                #Search TXT files
                if ext == ".pdf":
                    try:
                        doc = fitz.open(full_filename)
                        page = doc[0]
                        text_instances = page.searchFor(searchInput)
                        if text_instances:
                            print(full_filename)
                            fileList.insert(END, full_filename)
                            break
                    except:
                        pass

                elif ext == ".txt":
                    try:
                        fo = open(full_filename, 'r', encoding='utf-8', errors='ignore')

                        if searchInput in fo.read():
                            print(full_filename)
                            fileList.insert(END, full_filename)
                            break
                    except:
                        pass

                elif ext == ".docx":
                    try:
                        print(full_filename)
                        document = Document(full_filename)

                        for para in document.paragraphs:
                            if searchInput in para.text:
                                fileList.insert(END, full_filename)
                                break
                    except:
                        pass
                elif ext == ".xlsx":
                    try:
                        print(full_filename)
                        wbk = xlsxwriter.Workbook(full_filename)
                        wks = wbk.add_worksheet()

                        myPath = full_filename
                        for sh in xlrd.open_workbook(myPath).sheets():
                            for row in range(sh.nrows):
                                for col in range(sh.ncols):
                                    myCell = sh.cell(row, col)
                                    if myCell.value == searchInput:
                                        fileList.insert(END, full_filename)
                                        break
                                break
                    except:
                        pass
                elif ext == ".pptx":
                    try:
                        print(full_filename)
                        prs = Presentation(full_filename)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    if searchInput in shape.text:
                                        fileList.insert(END, full_filename)
                                    break
                            break
                    except:
                        pass

    except PermissionError:
        pass

# Search User's HDD. Send the information of HDD and filetypes.
def searchInfo():
    searchInput = searchText.get()
    HDD_List = []

    for drive in win32api.GetLogicalDriveStrings().split('\000')[:-1]:
        # insert the name of HDDs to the ListBox to display them for User to see.
        onlyHDDname = drive.split(":")
        desiredpath.delete(0, END)
        HDD_List.append(onlyHDDname[0])

        for i in HDD_List[0:]:
            desiredpath.insert(END, i + " ")

        # Specify the file_type.
        print("Finding files")
        searchDir(drive, searchInput)

# Execute Search Button
def executeSearchButton():
    searchInfo()

window = Tk()
C = Canvas(window, bg="blue", height=250)
backgroundimage = PhotoImage(file="image.png")
backgroundlabel = Label(window)
backgroundlabel.pack()

frame = Frame(window)
frame2 = Frame(window)

bOpen = tkinter.Button(window, text="open")
bClear = tkinter.Button(window, text="Clear")
bClearTwo = tkinter.Button(window, text="Clear")
btnSearch = tkinter.Button(window, text="search")
labelsearchtext = tkinter.Label(window, text="Enter text")
labledesiredpath = tkinter.Label(window, text="Searched HDD")
openedFileName = tkinter.Label(window, text="Opened Files")
searchedFileName = tkinter.Label(window, text="Searched Files")

fileList = Listbox(frame, width=50, height=17)
fileList.pack(side='left', fill='y')
scrollbar = Scrollbar(frame, orient=VERTICAL)
scrollbar.pack(side="right", fill="y")
scrollbar2 = Scrollbar(frame2, orient=VERTICAL)
scrollbar2.pack(side="right", fill="y")

openedFileList = Listbox(frame2, width=50)
openedFileList.pack(side='left', fill='y')
fileList.pack(side='left', fill='y')
desiredpath = tkinter.Listbox(window, width=10, height=5)
searchText = tkinter.Entry(window, width=16)

#fileList.bind('<<ListboxSelect>>', openFile)
#fileList.bind('<Double-Button>', buttonOpen)
openedFileList.bind('<Double-Button>', buttonOpen)
# initialise and run main loop
init(window)
mainloop()

